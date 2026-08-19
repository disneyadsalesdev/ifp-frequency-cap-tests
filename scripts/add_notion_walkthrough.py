"""Insert 99fold Notion Cursor walkthrough screenshots into the training deck."""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Reuse deck styling from modify_training_pptx
sys.path.insert(0, str(Path(__file__).resolve().parent))
from modify_training_pptx import (  # noqa: E402
    OUTPUT,
    SOURCE,
    add_header_bar,
    add_title_textbox,
    resolve_input_path,
)

NOTION_DIR = Path(__file__).resolve().parent / "pptx_assets" / "notion_cursor_walkthrough"
MANIFEST = NOTION_DIR / "manifest.json"
NOTION_URL = (
    "https://99fold.notion.site/Understanding-Cursor-A-Graphical-Walkthrough-of-the-Cursor-IDE-"
    "2ecb74f332388084b424f411bbbf719b"
)

WALKTHROUGH_SLIDES = [
    {
        "title": "Understanding Cursor — Graphical Walkthrough",
        "image": "walkthrough_02.png",
        "bullets": [
            "Adapted from the 99fold Notion guide by Matt Malishev (see speaker notes for link)",
            "Cursor combines file explorer, code editor, terminal, preview, and AI in one window",
            "Goal: less context switching — everything you need to ship work in one place",
        ],
    },
    {
        "title": "Step 1 — Use the File Explorer",
        "image": "walkthrough_03.png",
        "bullets": [
            "Browse and organise project files in a tree view on the left",
            "Expand folders to see structure at a glance",
            "All files are local on your machine via the Cursor desktop app",
        ],
    },
    {
        "title": "Step 2 — Work Inside the Code Editor",
        "image": "walkthrough_04.png",
        "bullets": [
            "Write and edit code across multiple tabs at the same time",
            "Switch between files using the tab bar at the top",
            "View live previews as you make changes (when your project supports it)",
        ],
    },
    {
        "title": "Step 3 — Run Commands in the Terminal",
        "image": "walkthrough_05.png",
        "bullets": [
            "Build, test, and launch apps or servers from the integrated terminal",
            "View logs and error messages where you work",
            "Run commands manually or let Cursor's AI execute them for you",
        ],
    },
    {
        "title": "Step 4 — Fix Code with the AI Assistant",
        "image": "walkthrough_06.png",
        "bullets": [
            "Chat with AI to resolve errors in real time",
            "Get explanations that reference your actual files and structure",
            "Review changes — Keep All to apply or Undo All to revert",
        ],
    },
]


def slide_title(slide) -> str:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip().split("\n")[0]
    return ""


def already_inserted(prs) -> bool:
    for slide in prs.slides:
        title = slide_title(slide)
        if "Step 1 — Use the File Explorer" in title or "Graphical Walkthrough" in title:
            return True
    return False


def find_insert_index(prs) -> int:
    """Insert after 'The Cursor Window' if present, else after 'Key Features'."""
    for i, slide in enumerate(prs.slides):
        title = slide_title(slide)
        if title.startswith("The Cursor Window"):
            return i + 1
    for i, slide in enumerate(prs.slides):
        title = slide_title(slide)
        if title.startswith("Key Features"):
            return i + 1
    return 6


def insert_slide_at(prs, index: int):
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    sld_id_lst = prs.slides._sldIdLst
    new_el = sld_id_lst[-1]
    sld_id_lst.remove(new_el)
    sld_id_lst.insert(index, new_el)
    return prs.slides[index]


def add_screenshot_slide(prs, index: int, title: str, bullets: list[str], image_path: Path) -> None:
    slide = insert_slide_at(prs, index)
    add_header_bar(slide, prs.slide_width)
    add_title_textbox(slide, title)

    img_left = Inches(0.45)
    img_top = Inches(1.35)
    img_width = Inches(5.6)
    img_height = Inches(4.55)
    slide.shapes.add_picture(str(image_path), img_left, img_top, img_width, img_height)

    body = slide.shapes.add_textbox(Inches(6.25), Inches(1.35), Inches(3.2), Inches(4.55))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(16)
        para.space_after = Pt(10)

    notes = slide.notes_slide.notes_text_frame
    notes.text = f"Source: 99fold Notion guide — {NOTION_URL}"


def apply_notion_walkthrough(prs) -> list[str]:
    if not MANIFEST.exists():
        raise FileNotFoundError(
            f"Run parse_notion_walkthrough.py first. Missing {MANIFEST}"
        )
    if already_inserted(prs):
        return ["Walkthrough slides already present — skipped"]

    log: list[str] = []
    insert_at = find_insert_index(prs)
    for offset, spec in enumerate(WALKTHROUGH_SLIDES):
        img_path = NOTION_DIR / spec["image"]
        if not img_path.exists():
            log.append(f"SKIP missing image: {img_path.name}")
            continue
        add_screenshot_slide(prs, insert_at + offset, spec["title"], spec["bullets"], img_path)
        log.append(f"Inserted at slide {insert_at + offset + 1}: {spec['title']}")

    return log


def main() -> int:
    input_path = OUTPUT if OUTPUT.exists() else resolve_input_path()
    if not MANIFEST.exists():
        print("Downloading Notion content first...")
        import parse_notion_walkthrough  # noqa: WPS433

        parse_notion_walkthrough.main()

    prs = Presentation(str(input_path))
    log = apply_notion_walkthrough(prs)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")
    for line in log:
        print(f"  • {line}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
