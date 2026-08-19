"""
Modify Cursor Training Guide PowerPoint:
- Remove section divider slides (SECTION 1, SECTION 2, …)
- Add/update Section 10 Claude & Copilot enterprise AI training slides
"""
from __future__ import annotations

import re
import sys
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Paths — prefer updated deck when re-running so Section 10 can be refreshed in place
SOURCE = Path(r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11.pptx")
OUTPUT = Path(r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11_updated.pptx")
ASSETS_DIR = Path(__file__).resolve().parent / "pptx_assets"
HEADER_RGB = (0x1E, 0x40, 0xAF)

# Section title slides: 0-based index -> visual config
SECTION_VISUALS: dict[int, dict] = {
    3: {
        "section": 1,
        "title": "Learn Cursor First",
        "subtitle": "Foundation & Key Concepts",
        "colors": ((30, 64, 175), (59, 130, 246)),
        "icon": "📚",
        "shapes": ["book", "lightbulb"],
    },
    7: {
        "section": 2,
        "title": "Download & Install",
        "subtitle": "Get Cursor Running on Your Machine",
        "colors": ((5, 150, 105), (16, 185, 129)),
        "icon": "⬇️",
        "shapes": ["download", "install"],
    },
    12: {
        "section": 3,
        "title": "Day 1 Setup",
        "subtitle": "Folders, Files & MCP Configuration",
        "colors": ((109, 40, 217), (139, 92, 246)),
        "icon": "🔌",
        "shapes": ["folder", "plug"],
    },
    23: {
        "section": 4,
        "title": "Cursor Settings",
        "subtitle": "Configure Your Workspace",
        "colors": ((55, 65, 81), (107, 114, 128)),
        "icon": "⚙️",
        "shapes": ["gear"],
    },
    26: {
        "section": 5,
        "title": "Chat & Agents",
        "subtitle": "New Agent, Run & Allowlist Buttons",
        "colors": ((13, 148, 136), (20, 184, 166)),
        "icon": "💬",
        "shapes": ["chat"],
    },
    31: {
        "section": 6,
        "title": "New Hire Startup",
        "subtitle": "First Prompts in Order",
        "colors": ((234, 88, 12), (251, 146, 60)),
        "icon": "🚀",
        "shapes": ["rocket"],
    },
    35: {
        "section": 7,
        "title": "IFP Every Day",
        "subtitle": "Daily Forecast Workflow",
        "colors": ((37, 99, 235), (96, 165, 250)),
        "icon": "📊",
        "shapes": ["chart"],
    },
    38: {
        "section": 8,
        "title": "DevTools Lab",
        "subtitle": "Capture Portal JSON for API Testing",
        "colors": ((185, 28, 28), (248, 113, 113)),
        "icon": "🔍",
        "shapes": ["devtools"],
    },
    43: {
        "section": 9,
        "title": "MSTR Daily Report",
        "subtitle": "Automated Reporting Workflow",
        "colors": ((21, 128, 61), (74, 222, 128)),
        "icon": "📋",
        "shapes": ["report"],
    },
}

# Image placement (matches existing slide picture layout)
IMG_LEFT = 457200
IMG_TOP = 960120
IMG_WIDTH = 8229600
IMG_HEIGHT = 5200000

# Section 10 — enterprise Claude & Copilot training (from SharePoint outline)
# No section divider slide — content slides only.
CLAUDE_COPILOT_SLIDES = [
    {
        "title": "How to Download & Access the Claude AI Portal",
        "bullets": [
            "Step 1: Open your company IT portal or AI enablement page — ask your manager for the Claude request link",
            "Step 2: Submit an access request for Claude Enterprise (team name + use case); manager approval may be required",
            "Step 3: When approved, open the enterprise URL from IT — do not use consumer claude.ai unless IT explicitly allows it",
            "Step 4: Sign in with company SSO (Microsoft Entra) — same credentials as email and SharePoint",
            "Step 5: Bookmark the portal; install desktop app only if IT lists it as an approved method",
        ],
    },
    {
        "title": "How to Integrate Claude with PowerPoint & Excel",
        "bullets": [
            "Excel — Step 1: Select your table or export CSV → paste into Claude → ask to summarize, pivot, or flag anomalies",
            "Excel — Step 2: Copy Claude's output back into a new sheet or slide; verify formulas and totals manually",
            "PowerPoint — Step 1: Paste your bullet outline → ask for slide titles, content, and speaker notes in brand tone",
            "PowerPoint — Step 2: Copy results into slides; use Copilot in-app (Copilot icon) if you have an M365 license — no copy/paste needed",
            "Rule: Copilot for quick in-document edits; Claude for deep analysis, long docs, or cross-tool work",
        ],
    },
    {
        "title": "How to Use Claude to Create HTML (Instead of Excel)",
        "bullets": [
            "Step 1: Copy your data from Excel (CSV) or paste a table — remove any confidential or PII columns first",
            "Step 2: Prompt: \"Create a self-contained HTML page with a styled table and bar chart for this data\"",
            "Step 3: Add: \"Use embedded CSS, responsive layout, and print-friendly styling — single .html file\"",
            "Step 4: Copy Claude's output → save as ReportName.html → double-click to open in browser and verify numbers",
            "Step 5: Share via approved internal site or email a link; use Excel when downstream formulas are still needed",
        ],
    },
    {
        "title": "How to Use Claude to Fill Out Email Reports",
        "bullets": [
            "Step 1: Collect inputs — weekly metrics, bullet notes, or meeting summary (follow data privacy policy)",
            "Step 2: Prompt: \"Draft a professional status email: subject line, Wins, Blockers, Next week — under 200 words\"",
            "Step 3: For recurring reports, add context: audience (leadership vs team), tone, and required sections",
            "Step 4: Review every name, date, metric, and attachment reference before sending — AI can hallucinate",
            "Step 5: In Outlook with Copilot: click Draft with Copilot in compose; save working prompts in SharePoint",
        ],
    },
    {
        "title": "How to Set Up Copilot Agents (Teams, SharePoint, M365)",
        "bullets": [
            "Step 1: Ask IT whether you have Microsoft 365 Copilot and/or Copilot Studio — both may require separate licenses",
            "Step 2 — Teams: Open a chat or channel → Copilot icon → try \"Summarize this channel\" or \"Draft a reply\"",
            "Step 3 — SharePoint: Open your site in browser → open Copilot pane → ask questions about pages and files on that site",
            "Step 4 — Custom agents: Go to copilotstudio.microsoft.com → Create agent (see next slides for full walkthrough)",
            "Step 5: Pin agents you use often; restrict access to approved Entra groups; test before team rollout",
        ],
    },
    {
        "title": "How to Build an Ask-Bot in Copilot Studio (SharePoint Q&A)",
        "bullets": [
            "Goal: a Teams bot that answers \"Where is the template?\" and points to the right SharePoint folder",
            "Step 1: Send IT request (Copilot Studio license + approved SharePoint folders) — template in docs/sharepoint-chatbot/",
            "Step 2: copilotstudio.microsoft.com → Create → New agent → name it (e.g. ET Anthropic Help Bot)",
            "Step 3: Paste system instructions; enable generative answers; require Entra ID sign-in",
            "Step 4: Knowledge → Add SharePoint site + specific libraries (Training, Templates) — not entire tenant",
        ],
    },
    {
        "title": "Copilot Studio — Connect SharePoint & Test the Bot",
        "bullets": [
            "Step 1: Knowledge → Add knowledge → SharePoint → select twdc.sharepoint.com/sites/ETAnthropic",
            "Step 2: Choose folders only: Training, Templates, AI Guides — exclude drafts and archive",
            "Step 3: Wait 15–60 min for indexing; open Test pane and run sample questions from config/sample-topics.json",
            "Step 4: Fix gaps — add missing docs to SharePoint or expand knowledge sources; re-test until answers cite real files",
            "Step 5: Bot must say \"I don't know\" when content isn't indexed — never invent file names or URLs",
        ],
    },
    {
        "title": "Copilot Studio — Publish Ask-Bot to Teams",
        "bullets": [
            "Step 1: Click Publish in Copilot Studio; wait for success confirmation",
            "Step 2: Channels → Microsoft Teams → turn On → choose org or specific Entra ID group",
            "Step 3: Teams → Apps → search your agent → Add to team → pick channel (e.g. #ai-training)",
            "Step 4: Post team quick-start with example prompts: \"Where is the Cursor guide?\" / \"Claude portal access?\"",
            "Step 5: Re-publish after major doc updates; review Analytics monthly for unanswered questions",
        ],
    },
    {
        "title": "Enterprise AI Best Practices (All Tools)",
        "bullets": [
            "Be specific in prompts: audience, format, length (\"3 bullets, no jargon, under 200 words\")",
            "Verify every output — AI can hallucinate dates, metrics, and file paths; cross-check source systems",
            "Never paste passwords, API keys, PII, or unreleased financials; follow acceptable-use policy",
            "Use the right tool: Cursor for code, Copilot for M365, Claude for deep analysis, Studio bot for SharePoint FAQ",
            "Save working prompts as team templates in SharePoint; share learnings in your AI Champions channel",
        ],
    },
    {
        "title": "SharePoint Q&A Bot — Team Quick Start",
        "bullets": [
            "Open ET Anthropic Help Bot in Teams or @mention it in your channel",
            "Try: \"Where is the Cursor training guide?\" / \"How do I get Claude portal access?\"",
            "Try: \"What's the MSTR daily report workflow?\" / \"Which folder has IFP templates?\"",
            "Bot answers from indexed SharePoint only — verify links before sharing externally",
            "Can't find an answer? Browse the SharePoint site or ask in #ai-training",
        ],
    },
]


def _load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def _draw_gradient(draw: ImageDraw.ImageDraw, w: int, h: int, c1: tuple, c2: tuple) -> None:
    for y in range(h):
        ratio = y / max(h - 1, 1)
        r = int(c1[0] + (c2[0] - c1[0]) * ratio)
        g = int(c1[1] + (c2[1] - c1[1]) * ratio)
        b = int(c1[2] + (c2[2] - c1[2]) * ratio)
        draw.line([(0, y), (w, y)], fill=(r, g, b))


def _draw_decorative_shapes(draw: ImageDraw.ImageDraw, w: int, h: int, kind: str, accent: tuple) -> None:
    light = tuple(min(c + 60, 255) for c in accent)
    if kind in ("book", "chart", "report"):
        draw.rounded_rectangle([w * 0.08, h * 0.15, w * 0.35, h * 0.75], radius=20, fill=light, outline="white", width=4)
        for i in range(4):
            y = h * 0.22 + i * h * 0.12
            draw.line([(w * 0.12, y), (w * 0.31, y)], fill="white", width=3)
    elif kind in ("download", "install"):
        cx = w * 0.2
        draw.polygon([(cx, h * 0.2), (cx + 60, h * 0.45), (cx - 60, h * 0.45)], fill="white")
        draw.rectangle([cx - 25, h * 0.42, cx + 25, h * 0.72], fill="white")
        draw.rectangle([cx - 70, h * 0.72, cx + 70, h * 0.78], fill="white")
    elif kind in ("folder", "plug"):
        draw.rounded_rectangle([w * 0.08, h * 0.28, w * 0.38, h * 0.62], radius=12, fill=light)
        draw.polygon([(w * 0.08, h * 0.28), (w * 0.18, h * 0.28), (w * 0.22, h * 0.22), (w * 0.38, h * 0.22)], fill=light)
        draw.ellipse([w * 0.55, h * 0.35, w * 0.75, h * 0.55], fill="white")
        draw.rectangle([w * 0.62, h * 0.42, w * 0.68, h * 0.48], fill=accent)
    elif kind == "gear":
        cx, cy, r = w * 0.22, h * 0.45, 80
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill="white", outline=light, width=6)
        draw.ellipse([cx - r * 0.4, cy - r * 0.4, cx + r * 0.4, cy + r * 0.4], fill=accent)
    elif kind == "chat":
        draw.rounded_rectangle([w * 0.06, h * 0.18, w * 0.42, h * 0.48], radius=20, fill="white")
        draw.polygon([(w * 0.12, h * 0.48), (w * 0.16, h * 0.58), (w * 0.22, h * 0.48)], fill="white")
        draw.rounded_rectangle([w * 0.18, h * 0.38, w * 0.48, h * 0.62], radius=20, fill=light)
    elif kind == "rocket":
        cx = w * 0.22
        draw.polygon([(cx, h * 0.2), (cx + 35, h * 0.55), (cx - 35, h * 0.55)], fill="white")
        draw.polygon([(cx - 50, h * 0.55), (cx - 20, h * 0.55), (cx - 35, h * 0.72)], fill=light)
        draw.polygon([(cx + 50, h * 0.55), (cx + 20, h * 0.55), (cx + 35, h * 0.72)], fill=light)
    elif kind == "devtools":
        draw.rounded_rectangle([w * 0.06, h * 0.15, w * 0.48, h * 0.55], radius=8, fill="white", outline=light, width=3)
        draw.rectangle([w * 0.06, h * 0.15, w * 0.48, h * 0.22], fill=accent)
        for i, c in enumerate(["#", "{", "</>"]):
            draw.text((w * 0.1 + i * w * 0.12, h * 0.28), c, fill=accent, font=_load_font(36, True))


def create_section_image(config: dict, out_path: Path) -> Path:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h))
    draw = ImageDraw.Draw(img)
    c1, c2 = config["colors"]
    _draw_gradient(draw, w, h, c1, c2)

    for shape in config.get("shapes", []):
        _draw_decorative_shapes(draw, w, h, shape, c1)

    badge_font = _load_font(28, True)
    badge_text = f"SECTION {config['section']}"
    bbox = draw.textbbox((0, 0), badge_text, font=badge_font)
    bw = bbox[2] - bbox[0] + 40
    draw.rounded_rectangle([w * 0.55, h * 0.12, w * 0.55 + bw, h * 0.12 + 50], radius=8, fill=(0, 0, 0, 80))
    draw.text((w * 0.55 + 20, h * 0.12 + 10), badge_text, fill="white", font=badge_font)

    icon_font = _load_font(120)
    draw.text((w * 0.55, h * 0.22), config.get("icon", "✦"), fill="white", font=icon_font)

    title_font = _load_font(52, True)
    sub_font = _load_font(32)
    draw.text((w * 0.55, h * 0.42), config["title"], fill="white", font=title_font)
    draw.text((w * 0.55, h * 0.52), config["subtitle"], fill=(230, 230, 255), font=sub_font)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path, "PNG")
    return out_path


def add_header_bar(slide, slide_width: int) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, 777240)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*HEADER_RGB)
    shape.line.fill.background()


def add_title_textbox(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(457200, 109728, 8229600, 594360)
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(28)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
    if subtitle:
        p1 = tf.add_paragraph()
        p1.text = subtitle
        p1.font.size = Pt(18)
        p1.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)


def _set_text_frame_content(text_frame, lines: list[str], *, title_mode: bool = False) -> None:
    """Replace all paragraphs in a text frame with new content."""
    text_frame.clear()
    for i, line in enumerate(lines):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.text = line
        para.level = 0
        if title_mode and i == 0:
            para.font.size = Pt(28)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
        elif title_mode and i > 0:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
        else:
            para.font.size = Pt(18)
            para.space_after = Pt(12)


def _text_boxes(slide) -> list:
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def update_section_title_slide(slide, title: str, subtitle: str | None, image_path: Path) -> None:
    boxes = _text_boxes(slide)
    if boxes:
        lines = [title] + ([subtitle] if subtitle else [])
        _set_text_frame_content(boxes[0].text_frame, lines, title_mode=True)
    # Refresh section visual if present
    for shape in list(slide.shapes):
        if shape.shape_type == 13:  # PICTURE
            slide.shapes._spTree.remove(shape._element)
    add_section_image_to_slide(slide, image_path)


def update_bullet_slide(slide, title: str, bullets: list[str]) -> None:
    boxes = _text_boxes(slide)
    if len(boxes) >= 1:
        _set_text_frame_content(boxes[0].text_frame, [title], title_mode=True)
    if len(boxes) >= 2:
        _set_text_frame_content(boxes[1].text_frame, bullets)


def add_bullet_slide(prs, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    add_header_bar(slide, prs.slide_width)
    add_title_textbox(slide, title)

    body = slide.shapes.add_textbox(548640, 960120, 8046720, 5200000)
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(18)
        para.space_after = Pt(12)


def add_section_slide(prs, title: str, subtitle: str | None, image_path: Path) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_header_bar(slide, prs.slide_width)
    add_title_textbox(slide, title, subtitle)
    add_section_image_to_slide(slide, image_path)


def add_section_image_to_slide(slide, image_path: Path) -> None:
    slide.shapes.add_picture(str(image_path), IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)


def find_section10_start(prs) -> int | None:
    """Find first Section 10 content slide (not divider-only slides)."""
    patterns = [
        re.compile(r"^How to Download", re.IGNORECASE),
        re.compile(r"^Getting Access to Claude", re.IGNORECASE),
        re.compile(r"^SECTION\s+10\s*[—\-]", re.IGNORECASE),
    ]
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                first_line = shape.text.strip().split("\n")[0]
                if any(p.search(first_line) for p in patterns):
                    return i
    return None


SECTION_DIVIDER_PATTERN = re.compile(r"^SECTION\s+\d+\s*[—\-]", re.IGNORECASE)


def is_section_divider_slide(slide) -> bool:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            first_line = shape.text.strip().split("\n")[0]
            if SECTION_DIVIDER_PATTERN.match(first_line):
                return True
    return False


def delete_slide(prs, index: int) -> None:
    slides = prs.slides._sldIdLst
    slide_id = slides[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    slides.remove(slide_id)


def remove_section_dividers(prs) -> list[str]:
    """Remove all SECTION N divider slides (highest index first)."""
    removed: list[str] = []
    indices = [i for i, slide in enumerate(prs.slides) if is_section_divider_slide(slide)]
    for idx in reversed(indices):
        slide = prs.slides[idx]
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip().split("\n")[0]
                break
        delete_slide(prs, idx)
        removed.append(f"Removed slide {idx + 1}: {title}")
    return removed


def is_blank_slide(slide) -> bool:
    """True when slide has no text and no pictures."""
    has_text = any(
        hasattr(shape, "text") and shape.text.strip() for shape in slide.shapes
    )
    has_picture = any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes
    )
    return not has_text and not has_picture


def remove_blank_slides(prs) -> list[str]:
    """Remove empty slides (highest index first)."""
    removed: list[str] = []
    indices = [i for i, slide in enumerate(prs.slides) if is_blank_slide(slide)]
    for idx in reversed(indices):
        delete_slide(prs, idx)
        removed.append(f"Removed blank slide {idx + 1}")
    return removed


def cleanup_slides(prs) -> tuple[list[str], list[str]]:
    """Remove SECTION divider slides and blank slides."""
    section_removed = remove_section_dividers(prs)
    blank_removed = remove_blank_slides(prs)
    return section_removed, blank_removed


def resolve_input_path() -> Path:
    """Always rebuild from the original deck so mid-deck slides are never overwritten."""
    if SOURCE.exists():
        return SOURCE
    if OUTPUT.exists():
        return OUTPUT
    print(f"ERROR: Neither source nor updated file found:\n  {SOURCE}\n  {OUTPUT}")
    sys.exit(1)


def apply_section_visuals(prs) -> list[str]:
    added_images: list[str] = []
    for slide_idx, config in SECTION_VISUALS.items():
        if slide_idx >= len(prs.slides):
            continue
        img_path = ASSETS_DIR / f"section_{config['section']}.png"
        create_section_image(config, img_path)
        slide = prs.slides[slide_idx]
        # Only add image if slide doesn't already have one
        has_pic = any(s.shape_type == 13 for s in slide.shapes)
        if not has_pic:
            add_section_image_to_slide(slide, img_path)
        added_images.append(
            f"Slide {slide_idx + 1} — SECTION {config['section']}: {config['title']}"
        )
    return added_images


def apply_claude_copilot_slides(prs) -> tuple[list[str], int]:
    """Update existing Section 10 slides or append if missing."""
    section10_start = find_section10_start(prs)
    log: list[str] = []
    content_slides = [s for s in CLAUDE_COPILOT_SLIDES if not s.get("is_section")]

    if section10_start is not None:
        existing_count = len(prs.slides) - section10_start
        for offset, slide_data in enumerate(content_slides):
            idx = section10_start + offset
            if idx >= len(prs.slides):
                break
            slide = prs.slides[idx]
            update_bullet_slide(slide, slide_data["title"], slide_data["bullets"])
            log.append(f"[updated] {slide_data['title']} ({len(slide_data['bullets'])} bullets)")

        if existing_count < len(content_slides):
            for slide_data in content_slides[existing_count:]:
                add_bullet_slide(prs, slide_data["title"], slide_data["bullets"])
                log.append(f"[added] {slide_data['title']} ({len(slide_data['bullets'])} bullets)")
    else:
        for slide_data in content_slides:
            add_bullet_slide(prs, slide_data["title"], slide_data["bullets"])
            log.append(f"[added] {slide_data['title']} ({len(slide_data['bullets'])} bullets)")

    start = (section10_start or len(prs.slides) - len(content_slides)) + 1
    return log, start


def main() -> None:
    input_path = resolve_input_path()
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(input_path))

    removed_sections, removed_blanks = cleanup_slides(prs)

    # Save and reload to avoid duplicate slide-part names after delete + append
    temp_path = OUTPUT.with_suffix(".temp.pptx")
    prs.save(str(temp_path))
    prs = Presentation(str(temp_path))

    notion_log: list[str] = []
    try:
        from add_notion_walkthrough import apply_notion_walkthrough

        notion_log = apply_notion_walkthrough(prs)
    except Exception as exc:
        notion_log = [f"Notion walkthrough skipped: {exc}"]

    claude_log, section10_first_slide = apply_claude_copilot_slides(prs)

    prs.save(str(OUTPUT))
    if SOURCE != OUTPUT:
        prs.save(str(SOURCE))
    if temp_path.exists():
        temp_path.unlink()
    print(f"Saved: {OUTPUT}")
    if SOURCE != OUTPUT:
        print(f"Saved: {SOURCE}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Section 10 content starts at slide: {section10_first_slide}")
    print("\n=== SECTION DIVIDERS REMOVED ===")
    for item in removed_sections:
        print(f"  • {item}")
    if not removed_sections:
        print("  • (none — already removed)")
    print("\n=== BLANK SLIDES REMOVED ===")
    for item in removed_blanks:
        print(f"  • {item}")
    if not removed_blanks:
        print("  • (none)")
    print("\n=== SECTION 10 — CLAUDE/COPILOT SLIDES ===")
    for item in claude_log:
        print(f"  • {item}")
    print("\n=== NOTION WALKTHROUGH (99fold) ===")
    for item in notion_log:
        print(f"  • {item}")
    print("\n=== SLIDE TITLES (Section 10) ===")
    content_count = len([s for s in CLAUDE_COPILOT_SLIDES if not s.get("is_section")])
    start_idx = section10_first_slide - 1
    for i in range(start_idx, min(start_idx + content_count, len(prs.slides))):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title_line = shape.text.strip().split("\n")[0]
                print(f"  Slide {i + 1}: {title_line}")
                break


if __name__ == "__main__":
    main()
