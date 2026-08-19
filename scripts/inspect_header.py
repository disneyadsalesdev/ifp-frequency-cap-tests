"""Check header rectangle styling."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR

path = r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11.pptx"
prs = Presentation(path)

slide = prs.slides[3]
for shape in slide.shapes:
    print(f"{shape.name}: fill={shape.fill.type}")
    if shape.fill.type is not None:
        try:
            print(f"  fore_color: {shape.fill.fore_color.rgb}")
        except Exception:
            try:
                print(f"  theme: {shape.fill.fore_color.theme_color}")
            except Exception:
                print("  color unknown")
