"""Get font sizes from content slides."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Pt

path = r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11.pptx"
prs = Presentation(path)

for idx in [1, 4, 8]:
    slide = prs.slides[idx]
    print(f"\n=== Slide {idx+1} ===")
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    f = run.font
                    rgb = f.color.rgb if f.color and f.color.type else "inherit"
                    print(f"  {run.text[:40]!r} -> size={f.size}, bold={f.bold}, rgb={rgb}")
