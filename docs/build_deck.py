"""Build Cursor 101 training PowerPoint deck (plain, simple language)."""
from __future__ import annotations

import subprocess
import sys
from pathlib import Path

DOCS_DIR = Path(__file__).parent
IMAGES_DIR = DOCS_DIR / "images"
OUTPUT = DOCS_DIR / "training-guide.pptx"

NAVY = (15, 23, 42)
BLUE = (30, 64, 175)
BLUE_LIGHT = (239, 246, 255)
WHITE = (255, 255, 255)
GRAY = (100, 116, 139)
DARK = (30, 41, 59)
CODE_BG = (241, 245, 249)
AMBER_BG = (255, 251, 235)
GREEN_BG = (240, 253, 244)


def ensure_pptx() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])


def rgb(color: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*color)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_bar(slide, title: str, prs, subtitle: str = "") -> None:
    from pptx.util import Inches, Pt

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(BLUE)
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(9), Inches(0.65))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26 if subtitle else 28)
    p.font.bold = True
    p.font.color.rgb = rgb(WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = rgb((191, 219, 254))


def add_bullets(slide, items: list[str], top: float = 1.05, size: int = 16, width: float = 8.8):
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(width), Inches(6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(DARK)
        p.space_after = Pt(5 if size <= 15 else 8)


def add_note(slide, text: str, top: float, kind: str = "tip"):
    from pptx.util import Inches, Pt

    bg = AMBER_BG if kind == "tip" else GREEN_BG
    border = (217, 119, 6) if kind == "tip" else (22, 163, 74)
    text_color = (146, 64, 14) if kind == "tip" else (21, 128, 61)

    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(top), Inches(8.8), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(bg)
    shape.line.color.rgb = rgb(border)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top + 0.1), Inches(8.5), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = rgb(text_color)


def add_code_block(slide, text: str, top: float, font_size: int = 12, max_height: float = 3.6):
    from pptx.util import Inches, Pt

    lines = text.strip().split("\n")
    line_h = 0.22 if font_size <= 10 else 0.32
    height = min(line_h * len(lines) + 0.28, max_height)
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(top), Inches(8.8), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(CODE_BG)
    shape.line.color.rgb = rgb(BLUE)

    box = slide.shapes.add_textbox(Inches(0.75), Inches(top + 0.1), Inches(8.5), Inches(height - 0.15))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb(DARK)


def add_table(
    slide,
    headers: tuple[str, ...],
    rows: list[tuple[str, ...]],
    top: float,
    row_h: float = 0.42,
    width: float = 8.8,
):
    from pptx.util import Inches, Pt

    ncols = len(headers)
    n = len(rows) + 1
    table = slide.shapes.add_table(n, ncols, Inches(0.6), Inches(top), Inches(width), Inches(row_h * n)).table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(BLUE)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11 if ncols <= 2 else 10)
            p.font.color.rgb = rgb(WHITE)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(BLUE_LIGHT)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10 if ncols > 2 else 11)
                p.font.color.rgb = rgb(DARK)


def add_flow_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    subtitle: str = "",
    fill: tuple[int, int, int] = BLUE_LIGHT,
) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(BLUE)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = rgb(DARK)
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.color.rgb = rgb(GRAY)
        p2.alignment = PP_ALIGN.CENTER


def add_arrow_label(slide, left: float, top: float, text: str = "→") -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(0.35), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = rgb(BLUE)
    p.alignment = PP_ALIGN.CENTER


def add_mcp_pipeline_diagram(slide, top: float = 1.15) -> None:
    """Horizontal flow: You → Cursor → MCP Server → External API."""
    boxes = [
        (0.55, "You", "Ask in chat"),
        (2.75, "Cursor", "Picks the helper"),
        (4.95, "MCP Server", "Go-between"),
        (7.15, "IFP API", "Gets the numbers"),
    ]
    box_w, box_h = 1.85, 0.95
    for i, (left, title, subtitle) in enumerate(boxes):
        add_flow_box(slide, left, top, box_w, box_h, title, subtitle)
        if i < len(boxes) - 1:
            add_arrow_label(slide, left + box_w + 0.05, top + 0.28)


def add_mcp_without_with_diagram(slide, top: float = 1.15) -> None:
    """Two-row comparison: without MCP vs with MCP."""
    from pptx.util import Inches, Pt

    label = slide.shapes.add_textbox(Inches(0.6), Inches(top - 0.05), Inches(2.0), Inches(0.35))
    p = label.text_frame.paragraphs[0]
    p.text = "Without MCP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = rgb(DARK)

    add_flow_box(slide, 0.6, top + 0.35, 1.5, 0.75, "You", "Ask in chat")
    add_arrow_label(slide, 2.15, top + 0.55)
    add_flow_box(slide, 2.55, top + 0.35, 1.7, 0.75, "Cursor", "Runs script")
    add_arrow_label(slide, 4.3, top + 0.55)
    add_flow_box(slide, 4.7, top + 0.35, 2.0, 0.75, "Terminal / Script", "run-forecasts.ps1")
    add_arrow_label(slide, 6.75, top + 0.55)
    add_flow_box(slide, 7.15, top + 0.35, 1.8, 0.75, "IFP API", "Forecast data")

    top2 = top + 1.65
    label2 = slide.shapes.add_textbox(Inches(0.6), Inches(top2 - 0.05), Inches(2.0), Inches(0.35))
    p = label2.text_frame.paragraphs[0]
    p.text = "With MCP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = rgb(BLUE)

    add_flow_box(slide, 0.6, top2 + 0.35, 1.5, 0.75, "You", "Ask in chat", fill=GREEN_BG)
    add_arrow_label(slide, 2.15, top2 + 0.55)
    add_flow_box(slide, 2.55, top2 + 0.35, 1.7, 0.75, "Cursor", "Calls MCP tool", fill=GREEN_BG)
    add_arrow_label(slide, 4.3, top2 + 0.55)
    add_flow_box(slide, 4.7, top2 + 0.35, 2.0, 0.75, "MCP Server", "run_forecast", fill=GREEN_BG)
    add_arrow_label(slide, 6.75, top2 + 0.55)
    add_flow_box(slide, 7.15, top2 + 0.35, 1.8, 0.75, "IFP API", "Forecast data", fill=GREEN_BG)


def add_mcp_setup_steps_diagram(slide, top: float = 1.1) -> None:
    """Vertical setup checklist as numbered boxes."""
    steps = [
        ("1", "Create MCP server", "Small program with tools\ne.g. ifp-mcp-server/server.py"),
        ("2", "Register in mcp.json", "Tell Cursor how to start it\n~/.cursor/mcp.json"),
        ("3", "Restart Cursor", "Reload so MCP is picked up"),
        ("4", "Ask in chat", 'e.g. "Run a forecast for Hulu US"'),
    ]
    y = top
    for num, title, detail in steps:
        from pptx.util import Inches, Pt

        badge = slide.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(BLUE)
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = num
        bp.font.size = Pt(16)
        bp.font.bold = True
        bp.font.color.rgb = rgb(WHITE)

        add_flow_box(slide, 1.2, y, 3.2, 0.85, title, detail.replace("\n", " — "))
        if num != "4":
            add_arrow_label(slide, 2.1, y + 0.88, "↓")
        y += 1.05


def add_image_or_placeholder(
    slide,
    filename: str,
    left: float,
    top: float,
    width: float,
    height: float,
    caption: str = "",
) -> None:
    """Embed a screenshot from docs/images/ or draw a labeled placeholder."""
    from pptx.util import Inches, Pt

    path = IMAGES_DIR / filename
    if path.is_file():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(CODE_BG)
        shape.line.color.rgb = rgb(GRAY)
        shape.line.dash_style = 2
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Screenshot: {filename}]\nDrop PNG in docs\\images\\ and rebuild deck"
        p.font.size = Pt(11)
        p.font.color.rgb = rgb(GRAY)
    if caption:
        cap = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.05), Inches(width), Inches(0.35))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(10)
        cp.font.color.rgb = rgb(GRAY)


def add_gfg_image_slide(
    prs,
    title: str,
    filename: str,
    subtitle: str = "",
    credit: str = "Visual reference: GeeksforGeeks Cursor AI guide",
) -> None:
    """Full-width infographic slide from docs/images/ (e.g. gfg-*.png)."""
    from pptx.util import Inches, Pt

    slide = blank_slide(prs)
    add_title_bar(slide, title, prs, subtitle)
    add_image_or_placeholder(slide, filename, 0.55, 1.05, 8.9, 5.35, "")
    if credit:
        foot = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(8.9), Inches(0.45))
        p = foot.text_frame.paragraphs[0]
        p.text = credit
        p.font.size = Pt(9)
        p.font.color.rgb = rgb(GRAY)


def add_cursor_ui_mockup(slide, top: float = 1.05) -> None:
    """Wireframe of Cursor: explorer | editor | chat + terminal strip."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    outer = slide.shapes.add_shape(1, Inches(0.55), Inches(top), Inches(8.9), Inches(4.35))
    outer.fill.solid()
    outer.fill.fore_color.rgb = rgb((248, 250, 252))
    outer.line.color.rgb = rgb(BLUE)

    panels = [
        (0.65, 1.75, "Explorer", "Project files", (226, 232, 240)),
        (2.45, 4.2, "Editor", "Open file / diff view", WHITE),
        (6.7, 2.7, "Chat (Ctrl+I)", "Agent · Ask · Plan", BLUE_LIGHT),
    ]
    for left, width, title, sub, fill in panels:
        box = slide.shapes.add_shape(1, Inches(left), Inches(top + 0.15), Inches(width), Inches(3.05))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill)
        box.line.color.rgb = rgb(BLUE)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = rgb(DARK)
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(9)
        p2.font.color.rgb = rgb(GRAY)
        p2.alignment = PP_ALIGN.CENTER

    term = slide.shapes.add_shape(1, Inches(0.65), Inches(top + 3.3), Inches(8.7), Inches(0.85))
    term.fill.solid()
    term.fill.fore_color.rgb = rgb((30, 41, 59))
    term.line.color.rgb = rgb(BLUE)
    tp = term.text_frame.paragraphs[0]
    tp.text = "Terminal (Ctrl+`) — PowerShell, Python, git"
    tp.font.size = Pt(10)
    tp.font.color.rgb = rgb(WHITE)


def add_mcp_connection_types_diagram(slide, top: float = 1.15) -> None:
    """Local stdio MCP vs remote HTTP MCP (e.g. Snowflake)."""
    from pptx.util import Inches, Pt

    lbl = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(4.2), Inches(0.3))
    lbl.text_frame.paragraphs[0].text = "Type A — Local MCP (stdio)"
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.size = Pt(12)

    add_flow_box(slide, 0.6, top + 0.35, 1.6, 0.7, "Cursor", "Starts process")
    add_arrow_label(slide, 2.25, top + 0.5)
    add_flow_box(slide, 2.65, top + 0.35, 2.0, 0.7, "MCP server.py", "Runs on your PC")
    add_arrow_label(slide, 4.7, top + 0.5)
    add_flow_box(slide, 5.1, top + 0.35, 1.9, 0.7, "Internal API", "IFP forecast")

    top2 = top + 1.35
    lbl2 = slide.shapes.add_textbox(Inches(0.6), Inches(top2), Inches(4.5), Inches(0.3))
    lbl2.text_frame.paragraphs[0].text = "Type B — Remote MCP (HTTPS + token)"
    lbl2.text_frame.paragraphs[0].font.bold = True
    lbl2.text_frame.paragraphs[0].font.size = Pt(12)
    lbl2.text_frame.paragraphs[0].font.color.rgb = rgb(BLUE)

    add_flow_box(slide, 0.6, top2 + 0.35, 1.6, 0.7, "Cursor", "HTTP client", fill=GREEN_BG)
    add_arrow_label(slide, 2.25, top2 + 0.5)
    add_flow_box(slide, 2.65, top2 + 0.35, 2.4, 0.7, "Snowflake MCP URL", "Hosted by Snowflake", fill=GREEN_BG)
    add_arrow_label(slide, 5.1, top2 + 0.5)
    add_flow_box(slide, 5.5, top2 + 0.35, 1.5, 0.7, "Warehouse", "SQL / Cortex", fill=GREEN_BG)


def add_snowflake_mcp_setup_diagram(slide, top: float = 1.05) -> None:
    steps = [
        ("1", "Create PAT", "Snowsight → Settings → Programmatic Access Token"),
        ("2", "Get MCP URL", "account / database / schema / mcp-servers / name"),
        ("3", "Set env vars", "SNOWFLAKE_MCP_SERVER_URL, SNOWFLAKE_PAT_TOKEN"),
        ("4", "Register in mcp.json", "url + Authorization: Bearer ${...}"),
        ("5", "Restart Cursor", "Settings → MCP → green status"),
    ]
    y = top
    for num, title, detail in steps:
        from pptx.util import Inches, Pt

        badge = slide.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(BLUE)
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = num
        bp.font.size = Pt(16)
        bp.font.bold = True
        bp.font.color.rgb = rgb(WHITE)

        add_flow_box(slide, 1.2, y, 7.2, 0.72, title, detail)
        if num != "5":
            add_arrow_label(slide, 4.5, y + 0.74, "↓")
        y += 0.95


def add_agent_toolbelt_diagram(slide, top: float = 1.2) -> None:
    """What the Agent can use on your behalf (conceptual)."""
    tools = [
        (0.55, "Read / edit\nproject files"),
        (2.15, "Search\ncodebase"),
        (3.75, "Run\nterminal"),
        (5.35, "Call MCP\ntools"),
        (6.95, "Use @ context\n(files, docs)"),
    ]
    for left, label in tools:
        add_flow_box(slide, left, top, 1.45, 0.95, label.replace("\n", " "), "", fill=BLUE_LIGHT)
    add_flow_box(slide, 2.75, top + 1.25, 4.5, 0.75, "You approve diffs & commands", "Cursor asks before risky actions", fill=AMBER_BG)


def title_slide(prs, title: str, subtitle: str, tag: str = ""):
    from pptx.util import Inches, Pt

    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(NAVY)
    bg.line.fill.background()
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(BLUE)
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = rgb(WHITE)

    s = slide.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(8.8), Inches(1.2))
    p = s.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = rgb((148, 163, 184))

    if tag:
        g = slide.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(8.8), Inches(0.5))
        p = g.text_frame.paragraphs[0]
        p.text = tag
        p.font.size = Pt(13)
        p.font.color.rgb = rgb((148, 163, 184))


def build() -> Path:
    ensure_pptx()
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    import sys
    from pathlib import Path as _Path

    _docs = _Path(__file__).resolve().parent
    if str(_docs) not in sys.path:
        sys.path.insert(0, str(_docs))
    from deck_onboarding import build_full_onboarding

    build_full_onboarding(prs)

    # REFERENCE — appendix (optional)
    slide = blank_slide(prs)
    add_title_bar(slide, "Appendix", prs, "Optional — main training is Sections 1–10 above")
    add_bullets(
        slide,
        [
            "Glossary, team jobs, MCP detail, cheat-sheet phrases below",
            "Use Table of Contents slide to jump back to any section",
        ],
        size=16,
    )

    # Glossary (one slide)
    slide = blank_slide(prs)
    add_title_bar(slide, "Glossary", prs, "Words you will hear")
    add_table(
        slide,
        ("Word", "Simple meaning"),
        [
            ("Prompt", "What you type in chat"),
            ("Agent / Ask", "Agent acts (with OK); Ask only explains"),
            ("MCP", "Plug-in tools — IFP forecast, Snowflake, etc."),
            ("Diff", "Before/after on a file — Accept or Reject"),
            ("@ mention", "Tag a file so Cursor knows which one"),
            ("Skill / Rule", "Team cheat sheets Cursor should follow"),
        ],
        top=1.05,
        row_h=0.42,
    )

    # Team folders
    slide = blank_slide(prs)
    add_title_bar(slide, "Team Folders (If You Have Access)", prs, "Skip daily-report folder until after IFP practice")
    add_table(
        slide,
        ("Folder name", "What job it does"),
        [
            ("all-publishers-daily-report", "Makes the daily Excel report and email draft"),
            ("ifp-frequency-cap-tests", "Gets forecast numbers and cap tests"),
        ],
        top=1.2,
        row_h=0.6,
    )
    add_bullets(
        slide,
        [
            "",
            "Open it: File → Open Folder → pick the folder",
            "",
            "Tip: Only open the folder you need today.",
        ],
        top=2.8,
        size=16,
    )

    slide = blank_slide(prs)
    add_title_bar(slide, "Every Time — 5 Easy Steps", prs, "Same recipe every job")
    add_bullets(
        slide,
        [
            "Step 1 — Open the right project folder",
            "",
            "Step 2 — Press Ctrl+I to open chat",
            "",
            "Step 3 — New Chat (+) → pick Agent mode (when you want something done)",
            "",
            "Step 4 — Type what you want in normal words",
            "",
            "Step 5 — Run/Approve if asked → Accept or Reject file changes",
        ],
        size=18,
    )

    # 10d — Dashboards & outputs
    slide = blank_slide(prs)
    add_title_bar(slide, "Dashboards & Producing Outputs", prs)
    add_table(
        slide,
        ("Output", "How"),
        [
            ("Daily email + Excel", "MSTR schedule → Outlook → dailypublisherscript.py → draft"),
            ("IFP numbers", "Cursor MCP or run-forecasts.ps1 → output/results.json"),
            ("Cap PASS/FAIL", "validate-results.ps1 or validate_cap_test_matrix"),
            ("Live dashboard", "MSTR Library in browser; export or subscribe"),
            ("Meeting slides", "Export + summarize in Cursor; redact sensitive data"),
        ],
        top=1.05,
        row_h=0.42,
    )

    # 11 — Agent vs Ask
    slide = blank_slide(prs)
    add_title_bar(slide, "Three Helper Modes", prs, "Pick before you type")
    add_table(
        slide,
        ("Mode", "When to use it"),
        [
            ("Agent", "You want something DONE (run a job, fix a file)"),
            ("Ask", "You only want to LEARN (no changes — like reading a book)"),
            ("Plan", "Big job — make a to-do list first, then do it"),
        ],
        top=1.3,
        row_h=0.52,
    )
    add_bullets(
        slide,
        [
            "",
            "Most work days: use Agent.",
            "",
            "Use Ask when you are curious:",
            '  "Explain @dailypublisherscript.py like I am new"',
        ],
        top=3.2,
        size=16,
    )

    # 22 — Plan mode
    slide = blank_slide(prs)
    add_title_bar(slide, "Plan Mode = Make a List First", prs)
    add_bullets(
        slide,
        [
            "Use Plan when the job is big or touches many files.",
            "",
            "Like planning a school project before you start:",
            "  1. Tell Cursor your goal in Plan mode",
            "  2. It writes steps — nothing changes yet",
            "  3. You fix the list if needed",
            "  4. Then Agent does the steps one by one",
            "",
            "Skip Plan for tiny jobs like Run @dailypublisherscript.py",
        ],
        size=15,
    )

    # 23 — Agent toolbelt
    slide = blank_slide(prs)
    add_title_bar(slide, "What the Helper Can Reach", prs, "You always get to say yes or no")
    add_agent_toolbelt_diagram(slide, top=1.15)
    add_bullets(
        slide,
        [
            "",
            "Cursor will ask before it runs commands or changes files.",
            "Never type passwords or secret keys in chat.",
        ],
        top=2.55,
        size=14,
    )

    # 24 — Rules, Skills, Hooks
    slide = blank_slide(prs)
    add_title_bar(slide, "How We Teach Cursor Our Rules", prs)
    add_table(
        slide,
        ("Thing", "Like…"),
        [
            ("Rules", "Classroom rules on the wall — always follow these"),
            ("Skills", "A recipe card for a specific team job"),
            ("Hooks", "Automatic habit — e.g. tidy up when you save"),
            ("Your settings", "Your personal preferences"),
        ],
        top=1.1,
        row_h=0.48,
    )
    # 25 — @ mention types
    slide = blank_slide(prs)
    add_title_bar(slide, "The @ Symbol — Point at Stuff", prs)
    add_table(
        slide,
        ("Type @…", "Why"),
        [
            ("File / Folder", "This is the exact file I mean"),
            ("Highlighted text", "Look at these lines only"),
            ("Terminals", "Here is the error message — help!"),
            ("Docs / Web", "Read this page (if your team allows it)"),
        ],
        top=1.15,
        row_h=0.52,
    )

    # 26 — @ mentions (examples)
    slide = blank_slide(prs)
    add_title_bar(slide, "@ Examples", prs)
    add_bullets(
        slide,
        [
            "Type @ and start typing a file name.",
            "Pick from the list.",
            "",
            "That way Cursor does not grab the wrong file.",
            "",
            "Examples:",
        ],
        top=1.05,
        size=16,
    )
    add_code_block(
        slide,
        "Run @dailypublisherscript.py\n"
        "Update @config/base-request.json\n"
        "Explain @dailypublisherscript.py step by step",
        3.0,
    )

    # 27 — Three workflows overview
    slide = blank_slide(prs)
    add_title_bar(slide, "Three Jobs — Pick One", prs)
    add_table(
        slide,
        ("I want to…", "Which recipe"),
        [
            ("Make today's report email", "Job 1 — Daily Report"),
            ("Get forecast numbers", "Job 2 — IFP Forecast"),
            ("Ask a question or change a setting", "Job 3 — Ask for Help"),
        ],
        top=1.2,
        row_h=0.6,
    )
    # 14 — Workflow 1 intro
    slide = blank_slide(prs)
    add_title_bar(slide, "Job 1: Daily Report", prs, "In one sentence")
    add_bullets(
        slide,
        [
            "Cursor grabs today's files from email,",
            "fills in the Excel workbook, and opens an email draft.",
            "",
            "It does NOT send the email — YOU click Send when it looks right.",
            "",
            "Folder: all-publishers-daily-report",
            "",
            "Magic sentence to type:",
        ],
        top=1.05,
        size=16,
    )
    add_code_block(slide, "Run @dailypublisherscript.py", 4.2)

    # 15 — Workflow 1 steps
    slide = blank_slide(prs)
    add_title_bar(slide, "Job 1: What Happens Inside", prs, "Like an assembly line")
    add_bullets(
        slide,
        [
            "Find today's data email in Outlook and download attachments",
            "",
            "Copy daily rows into the big Excel file",
            "",
            "Copy quarterly rows into the same file",
            "",
            "Refresh the pivot tables (the summary tables)",
            "",
            "Build summaries and open an Outlook draft",
            "",
            "You read it, fix anything, then send.",
        ],
        size=15,
    )

    # 16 — Workflow 1 prerequisites
    slide = blank_slide(prs)
    add_title_bar(slide, "Job 1: Check These First", prs)
    add_bullets(
        slide,
        [
            "Before you start:",
            "  • Outlook can open",
            "  • Today's data email already arrived",
            "  • The master Excel file is on your computer",
            "",
            "If it breaks, paste the error into chat:",
            '  "The email was not found — what should I check?"',
            "",
            "Backup if chat is broken:",
        ],
        top=1.05,
        size=15,
    )
    add_code_block(
        slide,
        "cd C:\\Users\\syeda012\\projects\\all-publishers-daily-report\npy dailypublisherscript.py",
        4.4,
    )

    # 17 — Workflow 2 intro
    slide = blank_slide(prs)
    add_title_bar(slide, "Job 2: Forecast Numbers", prs, "In one sentence")
    add_bullets(
        slide,
        [
            "Cursor asks the forecast computer for numbers and shows you the answer.",
            "",
            "You do NOT need to use the big website portal.",
            "",
            "Folder: ifp-frequency-cap-tests",
            "",
            "First: turn on Hulu VPN.",
            "",
            "Tell it: date, Hulu, US, city (optional), and cap (like 2 per hour).",
        ],
        top=1.05,
        size=15,
    )

    # 18 — Workflow 2 quick forecast
    slide = blank_slide(prs)
    add_title_bar(slide, "Job 2: One Quick Question", prs)
    add_bullets(slide, ["Type something like:"], top=1.05, size=16)
    add_code_block(slide, "Run a forecast for 7/10 Hulu US, Baton Rouge, 2 per hour", 1.5)
    add_bullets(
        slide,
        [
            "",
            "You get three numbers:",
            "",
            "  Capacity  — how much ads could show in total",
            "  Available — how much is left after the cap rule",
            "  Ratio     — Available divided by Capacity (like a percent)",
            "",
            "Example:",
            "  Capacity: 1,314,840  |  Available: 105,187  |  Ratio: 8.0%",
        ],
        top=2.5,
        size=15,
    )

    # 19 — Workflow 2 batch
    slide = blank_slide(prs)
    add_title_bar(slide, "Job 2: Run ALL the Tests", prs, "Optional — many caps at once")
    add_bullets(
        slide,
        [
            "Instead of one question, run every test in the list.",
            "",
            "Type:",
        ],
        top=1.05,
        size=16,
    )
    add_code_block(slide, "Run the frequency cap test matrix", 2.0)
    add_bullets(
        slide,
        [
            "",
            "Answers save to output\\results.json.",
            "",
            "Then ask:",
            '  "Show available and capacity for all scenarios"',
            "",
            "Cap examples:",
            "  2 per DAY  |  1 per HOUR  |  3 per 30 MINUTE  |  4 per 30 DAY",
        ],
        top=2.8,
        size=15,
    )

    # 20 — Workflow 3
    slide = blank_slide(prs)
    add_title_bar(slide, "Job 3: Just Ask", prs, "No special run button — talk normally")
    add_table(
        slide,
        ("I want to…", "Say…"),
        [
            ("Understand a file", '"Explain @dailypublisherscript.py step by step"'),
            ("Change who we target", '"Add Disney Plus to publisher targeting"'),
            ("Change dates", '"Update dates to July 15–31"'),
            ("City name → number", '"What\'s the DMA code for Baton Rouge?"'),
            ("Fix a mistake", '"Why did the last command fail?"'),
            ("Learn the steps", '"Walk me through a frequency cap test"'),
        ],
        top=1.05,
        row_h=0.42,
    )

    # 21 — Reviewing changes
    slide = blank_slide(prs)
    add_title_bar(slide, "When Cursor Changes a File", prs, "Always look before you say yes")
    add_bullets(
        slide,
        [
            "You see red and green — like track changes in Word.",
            "",
            "Your choices:",
            "  • Accept — keep it",
            "  • Reject — put it back",
            "",
            "Cursor will not send email or save secrets by itself.",
            "",
            "Not sure? Reject and ask:",
            '  "Explain that change first"',
        ],
        size=16,
    )

    # 36 — What is MCP?
    slide = blank_slide(prs)
    add_title_bar(slide, "What Is MCP?", prs, "Extra tools Cursor can plug in — like USB for helpers")
    add_bullets(
        slide,
        [
            "MCP is a shared language so Cursor can use outside tools safely.",
            "",
            "Think of labeled buttons:",
            '  \"Get forecast\" or \"Ask Snowflake a question\".',
            "",
            "A small program (MCP server) offers those buttons.",
            "",
            "Our forecast buttons include:",
            "  • run_forecast — one forecast",
            "  • run_cap_test_matrix — all tests",
            "  • validate_cap_test_matrix — check pass/fail",
            "",
            "You talk normally. Cursor picks the button. The server does the work.",
        ],
        size=14,
    )

    # 37 — MCP vocabulary
    slide = blank_slide(prs)
    add_title_bar(slide, "MCP Words — Simple", prs)
    add_table(
        slide,
        ("Word", "Simple meaning"),
        [
            ("MCP client", "Cursor — the app you use"),
            ("MCP server", "The go-between program"),
            ("Tool", "One job the server can do (like run_forecast)"),
            ("stdio", "Server runs on YOUR computer; Cursor talks to it locally"),
            ("HTTPS", "Server lives on the internet (like Snowflake)"),
            ("mcp.json", "A contact list telling Cursor how to call each server"),
        ],
        top=1.1,
        row_h=0.45,
    )

    # 38 — MCP connection types
    slide = blank_slide(prs)
    add_title_bar(slide, "Two Ways to Plug In MCP", prs, "On your laptop OR on the internet")
    add_mcp_connection_types_diagram(slide, top=1.05)

    # 39 — MCP pipeline diagram
    slide = blank_slide(prs)
    add_title_bar(slide, "How MCP Works (Story Version)", prs)
    add_mcp_pipeline_diagram(slide, top=1.25)
    add_bullets(
        slide,
        [
            "",
            "1. You ask: \"Run a forecast for Hulu US, 2 per hour\"",
            "2. Cursor picks the right MCP button (run_forecast)",
            "3. The go-between asks the forecast computer",
            "4. Numbers come back to chat",
            "",
            "You did not have to type scary commands yourself.",
        ],
        top=2.5,
        size=14,
    )

    # 24 — Without vs With MCP
    slide = blank_slide(prs)
    add_title_bar(slide, "With MCP vs Without MCP", prs)
    add_mcp_without_with_diagram(slide, top=1.05)

    # 25 — MCP user workflow
    slide = blank_slide(prs)
    add_title_bar(slide, "Using MCP (You)", prs, "If someone already set it up, you just chat")
    add_bullets(
        slide,
        [
            "Step 1 — Turn on Hulu VPN",
            "",
            "Step 2 — Open ifp-frequency-cap-tests folder",
            "",
            "Step 3 — Ctrl+I → Agent mode",
            "",
            "Step 4 — Type normal words, for example:",
        ],
        top=1.05,
        size=15,
    )
    add_code_block(
        slide,
        "Run a forecast for 7/10 Hulu US, Baton Rouge, 2 per hour\n\n"
        "Run the frequency cap test matrix\n\n"
        "Validate all cap cases",
        3.5,
    )
    add_bullets(
        slide,
        ["", "Step 5 — Read the results in chat. Done."],
        top=5.5,
        size=15,
    )

    # 26 — MCP setup overview
    slide = blank_slide(prs)
    add_title_bar(slide, "Setting Up MCP (Grown-Ups / IT)", prs, "One-time — like installing a new app button")
    add_mcp_setup_steps_diagram(slide, top=1.05)

    # 27 — MCP server pieces
    slide = blank_slide(prs)
    add_title_bar(slide, "MCP Has Two Parts", prs)
    add_table(
        slide,
        ("Part", "What it is"),
        [
            ("Helper program", "server.py — does the real work and calls forecast API"),
            ("Contact list", "mcp.json — tells Cursor how to start the helper"),
        ],
        top=1.1,
        row_h=0.55,
    )
    add_bullets(
        slide,
        [
            "",
            "Chat buttons: run_forecast | run_cap_test_matrix | validate_cap_test_matrix",
            "",
            "It reads settings from this project folder (dates, tests, etc.).",
        ],
        top=2.5,
        size=14,
    )

    # 28 — mcp.json example
    slide = blank_slide(prs)
    add_title_bar(slide, "mcp.json Example", prs, "Lives in C:\\Users\\<you>\\.cursor\\mcp.json")
    add_bullets(slide, ["Under mcpServers, add something like:"], top=1.05, size=14)
    add_code_block(
        slide,
        '{\n'
        '  "mcpServers": {\n'
        '    "ifp-forecast": {\n'
        '      "type": "stdio",\n'
        '      "command": "py",\n'
        '      "args": ["C:/.../ifp-mcp-server/server.py"],\n'
        '      "env": {\n'
        '        "IFP_API_URL": "http://inventory-forecasting-prod-dplus...",\n'
        '        "IFP_TESTS_ROOT": "C:/.../ifp-frequency-cap-tests"\n'
        '      }\n'
        '    }\n'
        '  }\n'
        '}',
        1.45,
    )
    add_bullets(
        slide,
        [
            "",
            "stdio = talk on your own computer",
            "command + args = how to start the helper",
            "env = paths and settings — never put passwords here",
            "",
            "Save → restart Cursor (or reload MCP in Settings)",
        ],
        top=4.0,
        size=13,
    )

    # 44 — Build your own MCP (official pattern + IFP)
    slide = blank_slide(prs)
    add_title_bar(slide, "Build Your Own MCP Server", prs, "Official guide → modelcontextprotocol.io/docs/develop/build-server")
    add_table(
        slide,
        ("Step", "Weather tutorial", "Our IFP server"),
        [
            ("1", "Create project folder", "projects/ifp-mcp-server"),
            ("2", "pip install mcp", "requirements.txt: mcp>=1.9.0"),
            ("3", "FastMCP + @mcp.tool()", "run_forecast, run_cap_test_matrix, …"),
            ("4", "mcp.run() at bottom", "stdio — Cursor launches via mcp.json"),
            ("5", "Add to mcp.json", "command py + absolute path to server.py"),
            ("6", "Reload Cursor", "Settings → MCP → green"),
        ],
        top=1.05,
        row_h=0.38,
    )
    add_bullets(
        slide,
        [
            "",
            "stdio servers: no print() to stdout — use logging (stderr) or MCP breaks.",
            "Tool docstrings (\"\"\" under each @mcp.tool) become the tool description in Cursor.",
        ],
        top=3.55,
        size=13,
    )

    # 45 — Snowflake MCP intro
    slide = blank_slide(prs)
    add_title_bar(slide, "Snowflake MCP (Another Plug-In)", prs, "Lives in the cloud — not on your laptop")
    add_bullets(
        slide,
        [
            "Snowflake hosts the helper. Cursor calls it over the internet.",
            "",
            "It might let you:",
            "  • Run SQL questions",
            "  • Search data with Cortex",
            "  • Use special team tools your admin turned on",
            "",
            "You get a secret token (PAT) — like a password. Never paste it in chat.",
            "Store it in Windows environment variables only.",
        ],
        size=14,
    )

    # 46 — Snowflake setup diagram
    slide = blank_slide(prs)
    add_title_bar(slide, "Snowflake Setup (Simple Steps)", prs)
    add_snowflake_mcp_setup_diagram(slide, top=1.0)

    # 47 — Snowflake mcp.json
    slide = blank_slide(prs)
    add_title_bar(slide, "Snowflake in mcp.json", prs, "URL and token come from env vars — not typed in the file")
    add_code_block(
        slide,
        '{\n'
        '  "mcpServers": {\n'
        '    "Snowflake": {\n'
        '      "url": "${SNOWFLAKE_MCP_SERVER_URL}",\n'
        '      "headers": {\n'
        '        "Authorization": "Bearer ${SNOWFLAKE_PAT_TOKEN}"\n'
        '      }\n'
        '    }\n'
        '  }\n'
        '}',
        1.15,
    )
    add_bullets(
        slide,
        [
            "",
            "Set SNOWFLAKE_MCP_SERVER_URL and SNOWFLAKE_PAT_TOKEN on your PC.",
            "Restart Cursor.",
            'Then chat: "Show tables in …" (if your team allows).',
        ],
        top=3.55,
        size=13,
    )

    # 48 — Snowflake vs IFP comparison
    slide = blank_slide(prs)
    add_title_bar(slide, "Forecast MCP vs Snowflake MCP", prs)
    add_table(
        slide,
        ("Question", "Forecast (IFP)", "Snowflake"),
        [
            ("Where it runs", "Your laptop", "Snowflake cloud"),
            ("How Cursor talks", "Local pipe (stdio)", "Internet + secret token"),
            ("Good for", "Ad forecast numbers", "SQL and warehouse data"),
            ("VPN", "Hulu VPN on", "Your company rules"),
        ],
        top=1.1,
        row_h=0.48,
    )

    # 49 — MCP setup in Cursor UI
    slide = blank_slide(prs)
    add_title_bar(slide, "Turn On MCP in Settings", prs, "Gear → search MCP — or edit mcp.json")
    add_bullets(
        slide,
        [
            "1. Ctrl+, (or gear) → search MCP",
            "2. Find ifp-forecast — green = ready, red = open logs",
            "3. Fix paths in C:\\Users\\<you>\\.cursor\\mcp.json if needed",
            "4. Ctrl+Shift+P → Developer: Reload Window",
            "5. VPN on before live IFP forecast in chat",
            "",
            "Green = ready. Red = wrong path, missing pip install mcp, or Python (use py).",
            "",
            "When agent calls an MCP tool, read the approval card → Run or Skip.",
        ],
        size=14,
    )

    # 30 — MCP troubleshooting
    slide = blank_slide(prs)
    add_title_bar(slide, "MCP Not Working?", prs)
    add_table(
        slide,
        ("Problem", "Try this"),
        [
            ("Not in Settings list", "Check mcp.json exists → Reload Window"),
            ("Server error / RED", "Settings → MCP → click server → read log"),
            ("Wrong path", "Use absolute paths with forward slashes in mcp.json"),
            ("Server hangs in terminal", "Normal for stdio — Cursor starts it, not you"),
            ("print() broke server", "Remove print(); use logging to stderr only"),
            ("Network error", "Turn on Hulu VPN for IFP"),
            ("Wrong numbers", "Say date, city, cap more clearly in chat"),
        ],
        top=1.05,
        row_h=0.36,
    )

    # 31 — Common problems
    slide = blank_slide(prs)
    add_title_bar(slide, "Oops — Fix It", prs)
    add_table(
        slide,
        ("Problem", "Try this"),
        [
            ('"Python was not found"', 'Tell Cursor: use py instead of python'),
            ("Network error", "Turn on Hulu VPN for forecast work"),
            ("Email not found", "Wait for today's data email in Outlook"),
            ("Wrong file changed", "Click Reject — be more specific next time"),
            ("Don't know what to type", "Use the cheat sheet slide"),
            ("City name vs number", 'Ask: "DMA code for [city]?"'),
        ],
        top=1.05,
        row_h=0.42,
    )

    # 32 — Prompts cheat sheet
    slide = blank_slide(prs)
    add_title_bar(slide, "Sentences You Can Copy", prs)
    add_code_block(
        slide,
        "PRACTICE LAB (ifp-frequency-cap-tests):\n"
        "Explain @reference/cap-ratio-expectations.json\n\n"
        "DMA code for Baton Rouge using @reference/dma-codes.json\n\n"
        "Explain @config/base-request.json (Ask only)\n\n"
        "Create practice-notes.md summarizing Steps 1–3\n\n"
        "Run a forecast for 7/10 Hulu US, Baton Rouge, 2 per hour (VPN)\n\n"
        "— OTHER TEAM JOBS —\n"
        "Run @dailypublisherscript.py\n\n"
        "Why did the last command fail?",
        1.05,
    )

    # 33 — Shortcuts
    slide = blank_slide(prs)
    add_title_bar(slide, "Keyboard Shortcuts", prs)
    add_table(
        slide,
        ("Do this", "Press"),
        [
            ("Open chat", "Ctrl+I"),
            ("Fix highlighted text", "Ctrl+K"),
            ("Open command box at bottom", "Ctrl+`"),
            ("Command menu", "Ctrl+Shift+P"),
            ("Send message fast", "Ctrl+Enter"),
            ("Use ghost suggestion", "Tab"),
        ],
        top=1.2,
        row_h=0.42,
    )
    add_bullets(
        slide,
        [
            "",
            "Stay safe:",
            "  • No passwords in chat",
            "  • Don't share forecast data outside the team",
            "  • Read changes before Accept",
        ],
        top=3.5,
        size=15,
    )

    # 35 — Closing
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(NAVY)
    bg.line.fill.background()
    from pptx.util import Pt

    t = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(8.8), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = "You got this!"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = rgb(WHITE)
    s = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(8.8), Inches(1.8))
    tf = s.text_frame
    tf.word_wrap = True
    for i, line in enumerate(
        [
            "Open folder → Ctrl+I → Agent → say what you need",
            "",
            "Cursor helps. You decide yes or no.",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = rgb((148, 163, 184))

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Build Cursor 101 training PowerPoint.")
    parser.add_argument(
        "--output",
        "-o",
        type=Path,
        default=OUTPUT,
        help="Path for .pptx (default: docs/training-guide.pptx)",
    )
    args = parser.parse_args()
    out = args.output.resolve()
    globals()["OUTPUT"] = out
    path = build()
    print(f"Created: {path} ({path.stat().st_size // 1024} KB)")
