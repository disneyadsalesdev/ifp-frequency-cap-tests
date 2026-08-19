"""Inspect text formatting on a sample slide."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11.pptx"
prs = Presentation(path)

slide = prs.slides[3]  # Section 1
for shape in slide.shapes:
    if not hasattr(shape, "text_frame"):
        continue
    print(f"Shape: {shape.name}")
    for pi, para in enumerate(shape.text_frame.paragraphs):
        print(f"  Para {pi}: level={para.level}, align={para.alignment}")
        for ri, run in enumerate(para.runs):
            f = run.font
            print(f"    Run: {run.text[:50]!r} size={f.size} bold={f.bold} color={f.color.rgb if f.color.type else None}")

# Check blank layout
print("\nLayouts:")
for layout in prs.slide_layouts:
    print(f"  {layout.name}")
