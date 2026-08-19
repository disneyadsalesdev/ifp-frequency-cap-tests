"""Verify modified PowerPoint."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

path = r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11_updated.pptx"
prs = Presentation(path)
print(f"Total slides: {len(prs.slides)}")

section_indices = [3, 7, 12, 23, 26, 31, 35, 38, 43, 74, 75, 76, 77, 78, 79, 80]
for idx in section_indices:
    if idx >= len(prs.slides):
        continue
    slide = prs.slides[idx]
    pics = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    texts = [s.text.strip().split("\n")[0][:60] for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
    print(f"Slide {idx+1}: pics={pics} | {texts[0] if texts else 'no text'}")

print("\nLast 7 slides:")
for i in range(max(0, len(prs.slides)-7), len(prs.slides)):
    slide = prs.slides[i]
    texts = [s.text.strip().split("\n")[0][:70] for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
    print(f"  {i+1}: {texts[0] if texts else '(empty)'}")
