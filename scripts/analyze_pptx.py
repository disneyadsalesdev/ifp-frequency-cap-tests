"""Analyze Cursor Training Guide PowerPoint - find section title slides."""
import re
import sys

sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11.pptx"
prs = Presentation(path)

section_pattern = re.compile(r"SECTION\s+\d+", re.IGNORECASE)

print(f"Total slides: {len(prs.slides)}\n")

section_slides = []

for i, slide in enumerate(prs.slides):
    all_text = []
    pics = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
        if hasattr(shape, "text") and shape.text.strip():
            all_text.append(shape.text.strip())

    combined = " | ".join(all_text)
    is_section = bool(section_pattern.search(combined))
    if is_section:
        section_slides.append(
            {
                "index": i,
                "slide_num": i + 1,
                "text": combined[:120],
                "pictures": pics,
                "shape_count": len(slide.shapes),
            }
        )

print("SECTION TITLE SLIDES:")
for s in section_slides:
    blank = s["pictures"] == 0 and s["shape_count"] <= 2
    print(
        f"  Slide {s['slide_num']}: {s['text'][:100]} | pics={s['pictures']} shapes={s['shape_count']} likely_blank={blank}"
    )

print("\nALL SLIDES (title only):")
for i, slide in enumerate(prs.slides):
    texts = []
    pics = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip().split("\n")[0][:80])
    title = texts[0] if texts else "(no text)"
    print(f"  {i+1:2d}. [{pics}pic] {title}")
