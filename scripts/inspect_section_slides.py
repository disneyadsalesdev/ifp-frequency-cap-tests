"""Inspect section title slide geometry."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

path = r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11.pptx"
prs = Presentation(path)

section_indices = [3, 7, 12, 23, 26, 31, 35, 38, 43]  # 0-based

for idx in section_indices:
    slide = prs.slides[idx]
    print(f"\n=== Slide {idx+1} ===")
    print(f"Layout: {slide.slide_layout.name}")
    for shape in slide.shapes:
        print(f"  {shape.name}: type={shape.shape_type}, left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        if hasattr(shape, "text") and shape.text.strip():
            print(f"    text: {shape.text.strip()[:100]!r}")

print(f"\nSlide dimensions: w={prs.slide_width}, h={prs.slide_height}")
print(f"In inches: w={prs.slide_width/914400:.2f}, h={prs.slide_height/914400:.2f}")
