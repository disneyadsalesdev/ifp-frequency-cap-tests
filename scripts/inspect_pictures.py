"""Inspect slides that already have pictures for layout reference."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r"c:\Users\syeda012\Downloads\Cursor_Training_Guide_8.11.pptx"
prs = Presentation(path)

for idx in [4, 6, 10, 11]:  # slides 5, 7, 11, 12
    slide = prs.slides[idx]
    print(f"\n=== Slide {idx+1} ===")
    for shape in slide.shapes:
        info = f"  {shape.name}: type={shape.shape_type}, left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}"
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            info += " [PICTURE]"
        if hasattr(shape, "text") and shape.text.strip():
            info += f" text={shape.text.strip()[:60]!r}"
        print(info)
